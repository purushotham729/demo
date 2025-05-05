import openai
import streamlit as st
from docx import Document
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

# ----- CONFIGURATION -----
OPENAI_API_KEY = 'your-openai-api-key-here'
MODEL = 'gpt-4-turbo'
MAX_TOKENS_PER_CHUNK = 2000
SUMMARY_MAX_TOKENS = 800

# ----- UTILITY FUNCTIONS -----
def read_docx(uploaded_file):
    doc = Document(uploaded_file)
    return [para.text.strip() for para in doc.paragraphs if para.text.strip()]

def read_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_runs.append(text)
    return text_runs

def chunk_paragraphs(paragraphs, max_tokens=MAX_TOKENS_PER_CHUNK):
    chunks, current_chunk, current_tokens = [], [], 0
    for para in paragraphs:
        tokens = len(para.split())
        if current_tokens + tokens > max_tokens:
            chunks.append("\n".join(current_chunk))
            current_chunk, current_tokens = [para], tokens
        else:
            current_chunk.append(para)
            current_tokens += tokens
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    return chunks

def summarize_chunk_safe(chunk, retry_count=3):
    for attempt in range(retry_count):
        try:
            response = openai.ChatCompletion.create(
                model=MODEL,
                api_key=OPENAI_API_KEY,
                messages=[
                    {"role": "system", "content": "Summarize the following document chunk clearly, retaining important requirements, features, and key points."},
                    {"role": "user", "content": chunk}
                ],
                temperature=0.2,
                max_tokens=1500
            )
            return response['choices'][0]['message']['content']
        except Exception as e:
            print(f"Error summarizing chunk (attempt {attempt+1}): {e}")
            time.sleep(2)
    return "[Error: Failed to summarize this chunk.]"

def summarize_document(paragraphs):
    chunks = chunk_paragraphs(paragraphs)
    summaries = [""] * len(chunks)
    progress_bar = st.progress(0)
    total = len(chunks)

    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {executor.submit(summarize_chunk_safe, chunk): i for i, chunk in enumerate(chunks)}
        completed = 0
        for future in as_completed(futures):
            i = futures[future]
            try:
                summaries[i] = future.result()
            except Exception as e:
                summaries[i] = "[Error]"
            completed += 1
            progress_bar.progress(completed / total)

    progress_bar.empty()
    return "\n\n".join(summaries)

def generate_new_frd(existing_brd_summary, existing_frd_summary, new_brd_summary):
    user_prompt = f"""
EXISTING BRD SUMMARY:
{existing_brd_summary}

EXISTING FRD SUMMARY:
{existing_frd_summary}

NEW BRD SUMMARY:
{new_brd_summary}

Please generate the NEW FRD.
"""
    system_prompt = (
        "You are an expert business analyst. "
        "You are given summarized versions of an existing BRD, FRD, and a new BRD. "
        "Your task is to create a NEW FRD based on the new BRD, maintaining structure and clarity of the existing FRD."
    )
    response = openai.ChatCompletion.create(
        model=MODEL,
        api_key=OPENAI_API_KEY,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
        max_tokens=3000
    )
    return response.choices[0].message.content

# ----- STREAMLIT UI -----
st.set_page_config(page_title="Auto FRD Generator", layout="centered")
st.title("📄 Auto FRD Generator")
st.markdown("Upload your existing BRD, FRD, and (optionally) new BRD file below to generate a new FRD.")

existing_brd_file = st.file_uploader("Upload Existing BRD (.docx or .pptx)", type=["docx", "pptx"])
existing_frd_file = st.file_uploader("Upload Existing FRD (.docx)", type="docx")
new_brd_file = st.file_uploader("Upload New BRD (.docx)", type="docx")

if st.button("Generate New FRD"):
    if not existing_brd_file or not existing_frd_file:
        st.error("❌ Please upload both Existing BRD and Existing FRD.")
    else:
        with st.spinner("Reading and summarizing documents..."):
            if existing_brd_file.name.endswith(".pptx"):
                paragraphs_brd = read_pptx(existing_brd_file)
            else:
                paragraphs_brd = read_docx(existing_brd_file)

            paragraphs_frd = read_docx(existing_frd_file)
            paragraphs_new_brd = read_docx(new_brd_file) if new_brd_file else []

            summary_brd = summarize_document(paragraphs_brd)
            summary_frd = summarize_document(paragraphs_frd)
            summary_new_brd = summarize_document(paragraphs_new_brd) if new_brd_file else "No new BRD provided."

        with st.spinner("Generating NEW FRD..."):
            new_frd_text = generate_new_frd(summary_brd, summary_frd, summary_new_brd)

        st.success("✅ FRD Generated Successfully!")
        st.download_button("📥 Download New FRD", new_frd_text, file_name="new_frd.txt", mime="text/plain")
        st.text_area("Preview of Generated FRD", new_frd_text, height=300)
